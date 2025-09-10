import os
from typing import List, Optional
from lxml import etree
from services.html_to_text_runs_service import (
    parse_html_text_to_text_runs as parse_inline_html_to_runs,
)

from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide
from pptx.text.text import _Paragraph, TextFrame, Font, _Run
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml.etree import fromstring, tostring
from PIL import Image, ImageDraw
import requests
import uuid
from pptx.oxml.xmlchemy import OxmlElement

from pptx.util import Pt
from pptx.dml.color import RGBColor

from models.pptx_models import (
    PptxAutoShapeBoxModel,
    PptxBoxShapeEnum,
    PptxConnectorModel,
    PptxFillModel,
    PptxFontModel,
    PptxParagraphModel,
    PptxPictureBoxModel,
    PptxPositionModel,
    PptxPresentationModel,
    PptxShadowModel,
    PptxSlideModel,
    PptxSpacingModel,
    PptxStrokeModel,
    PptxTextBoxModel,
    PptxTextRunModel,
)
from utils.download_helpers import download_files
from utils.get_env import get_app_data_directory_env
from utils.image_utils import (
    clip_image,
    create_circle_image,
    fit_image,
    invert_image,
    round_image_corners,
    set_image_opacity,
)
import uuid

BLANK_SLIDE_LAYOUT = 6


class PptxPresentationCreator:

    def __init__(self, ppt_model: PptxPresentationModel, temp_dir: str):
        self._temp_dir = temp_dir

        self._ppt_model = ppt_model
        self._slide_models = ppt_model.slides

        self._ppt = Presentation()
        self._ppt.slide_width = Pt(1280)
        self._ppt.slide_height = Pt(720)

    def get_sub_element(self, parent, tagname, **kwargs):
        """Helper method to create XML elements"""
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def _apply_smooth_border_radius(self, image: Image.Image, radii: List[int]) -> Image.Image:
        """Apply smooth anti-aliased border radius to image."""
        if not radii or all(r == 0 for r in radii):
            return image
            
        w, h = image.size
        max_radius = min(w // 2, h // 2)
        clamped_radii = [min(radius, max_radius) for radius in radii]
        
        # Use 4x supersampling for smooth edges
        scale = 4
        high_w, high_h = w * scale, h * scale
        high_radii = [r * scale for r in clamped_radii]
        
        # Scale up image
        high_image = image.resize((high_w, high_h), Image.LANCZOS)
        
        # Create high-res mask
        mask = Image.new("L", (high_w, high_h), 0)
        draw = ImageDraw.Draw(mask)
        
        # Draw main rectangle
        max_r = max(high_radii)
        draw.rectangle((max_r, 0, high_w - max_r, high_h), fill=255)
        draw.rectangle((0, max_r, high_w, high_h - max_r), fill=255)
        
        # Draw smooth corners using ellipses
        for i, radius in enumerate(high_radii):
            if radius > 0:
                if i == 0:  # top-left
                    draw.ellipse((0, 0, radius * 2, radius * 2), fill=255)
                elif i == 1:  # top-right
                    draw.ellipse((high_w - radius * 2, 0, high_w, radius * 2), fill=255)
                elif i == 2:  # bottom-right  
                    draw.ellipse((high_w - radius * 2, high_h - radius * 2, high_w, high_h), fill=255)
                elif i == 3:  # bottom-left
                    draw.ellipse((0, high_h - radius * 2, radius * 2, high_h), fill=255)
        
        # Apply mask
        high_image.putalpha(mask)
        
        # Scale back down with high quality
        return high_image.resize((w, h), Image.LANCZOS)
    
    def _create_smooth_circle_image(self, image: Image.Image) -> Image.Image:
        """Create smooth anti-aliased circle image."""
        w, h = image.size
        
        # Use 4x supersampling
        scale = 4
        high_w, high_h = w * scale, h * scale
        
        # Scale up image
        high_image = image.resize((high_w, high_h), Image.LANCZOS)
        
        # Create circular mask
        mask = Image.new("L", (high_w, high_h), 0)
        draw = ImageDraw.Draw(mask)
        
        # Calculate circle dimensions
        circle_size = min(high_w, high_h)
        radius = circle_size // 2
        center_x, center_y = high_w // 2, high_h // 2
        
        # Draw anti-aliased circle
        draw.ellipse((
            center_x - radius, center_y - radius,
            center_x + radius, center_y + radius
        ), fill=255)
        
        # Apply mask
        high_image.putalpha(mask)
        
        # Scale back down with high quality
        return high_image.resize((w, h), Image.LANCZOS)

    def _is_transparent_icon(self, image_path: str, local_file_path: str) -> bool:
        """Check if image is likely a transparent icon (PNG/SVG in icons directory)"""
        if not image_path:
            return False
            
        # Check if it's from icons directory
        is_from_icons = (
            "/static/icons/" in image_path or 
            "\\static\\icons\\" in image_path or
            (local_file_path and "icons" in local_file_path)
        )
        
        # Check file extension
        is_icon_format = image_path.lower().endswith(('.png', '.svg'))
        
        result = is_from_icons and is_icon_format
        print(f"is_transparent_icon check: image_path={image_path}, local_file_path={local_file_path}, is_from_icons={is_from_icons}, is_icon_format={is_icon_format}, result={result}")
        return result

    def _create_icon_with_background(self, local_file_path: str, picture_model: PptxPictureBoxModel) -> str:
        """Create icon with solid background to avoid transparency issues"""
        try:
            print(f"Creating icon with background: {local_file_path}")
            
            # Load the original icon
            icon = Image.open(local_file_path)
            print(f"Original icon mode: {icon.mode}, size: {icon.size}")
            
            # Convert to RGBA for processing
            if icon.mode != 'RGBA':
                icon = icon.convert('RGBA')
            
            # Get target size
            width = int(picture_model.position.width) if picture_model.position.width > 0 else icon.width
            height = int(picture_model.position.height) if picture_model.position.height > 0 else icon.height
            
            # Resize icon if needed
            if icon.size != (width, height):
                icon = icon.resize((width, height), Image.LANCZOS)
            
            # Create background - use white background by default
            # You can customize this color based on your needs
            background_color = (255, 255, 255, 255)  # White background
            
            # Create new image with solid background
            result_image = Image.new('RGBA', (width, height), background_color)
            
            # Paste the icon on top of the background
            # The icon's alpha channel will be used for blending
            result_image.paste(icon, (0, 0), icon)
            
            # Convert to RGB to remove alpha channel completely
            final_image = Image.new('RGB', (width, height), (255, 255, 255))
            final_image.paste(result_image, (0, 0), result_image)
            
            # Save processed image
            processed_path = os.path.join(self._temp_dir, f"icon_bg_{uuid.uuid4()}.png")
            final_image.save(processed_path, "PNG", quality=95, optimize=True)
            
            print(f"Created icon with background: {processed_path}")
            return processed_path
            
        except Exception as e:
            print(f"Failed to create icon with background {local_file_path}: {e}")
            # Return original file as fallback
            return local_file_path

    def _create_html_element_screenshot(self, image_path: str, picture_model: PptxPictureBoxModel) -> str:
        """Create screenshot of HTML element instead of processing transparent icon"""
        try:
            print(f"Creating screenshot for HTML element: {image_path}")
            
            # Extract the element information to create HTML for screenshot
            # For icons, we'll create a simple HTML div with the icon
            width = int(picture_model.position.width)
            height = int(picture_model.position.height)
            
            # Create HTML content for the icon element
            html_content = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                <style>
                    body {{
                        margin: 0;
                        padding: 0;
                        width: {width}px;
                        height: {height}px;
                        background: transparent;
                    }}
                    .icon-container {{
                        width: 100%;
                        height: 100%;
                        display: flex;
                        align-items: center;
                        justify-content: center;
                        background: transparent;
                    }}
                    .icon-img {{
                        width: 100%;
                        height: 100%;
                        object-fit: cover;
                        background: transparent;
                    }}
                </style>
            </head>
            <body>
                <div class="icon-container">
                    <img src="file://{local_file_path}" alt="icon" class="icon-img" />
                </div>
            </body>
            </html>
            """
            
            # Save HTML to temporary file
            html_file_path = os.path.join(self._temp_dir, f"icon_{uuid.uuid4()}.html")
            with open(html_file_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            # Use puppeteer to take screenshot
            screenshot_path = os.path.join(self._temp_dir, f"icon_screenshot_{uuid.uuid4()}.png")
            
            # Call puppeteer via subprocess
            import subprocess
            
            # First check if node and puppeteer are available
            try:
                # Check if node is available
                node_check = subprocess.run(["node", "--version"], capture_output=True, text=True, timeout=5)
                if node_check.returncode != 0:
                    print("Node.js not found, falling back to original file")
                    return local_file_path
                    
                # Check if puppeteer is available in nextjs directory
                nextjs_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "nextjs")
                if not os.path.exists(nextjs_dir):
                    print("Next.js directory not found, falling back to original file")
                    return local_file_path
                
                cmd = [
                    "node", 
                    "-e",
                    f"""
                    process.chdir('{nextjs_dir}');
                    const puppeteer = require('puppeteer');
                    (async () => {{
                        try {{
                            const browser = await puppeteer.launch({{
                                headless: true,
                                args: ['--no-sandbox', '--disable-setuid-sandbox', '--disable-dev-shm-usage']
                            }});
                            const page = await browser.newPage();
                            await page.setViewport({{width: {width}, height: {height}, deviceScaleFactor: 1}});
                            await page.goto('file://{html_file_path}', {{waitUntil: 'networkidle0', timeout: 10000}});
                            
                            // Ensure transparent background
                            await page.evaluate(() => {{
                                document.body.style.background = 'transparent';
                            }});
                            
                            await page.screenshot({{
                                path: '{screenshot_path}',
                                omitBackground: true,
                                type: 'png'
                            }});
                            await browser.close();
                            console.log('Screenshot completed successfully');
                        }} catch (error) {{
                            console.error('Screenshot error:', error.message);
                            process.exit(1);
                        }}
                    }})();
                    """
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=30, cwd=nextjs_dir)
                
                if result.returncode == 0 and os.path.exists(screenshot_path):
                    print(f"Successfully created screenshot: {screenshot_path}")
                    # Clean up temporary HTML file
                    if os.path.exists(html_file_path):
                        os.remove(html_file_path)
                    return screenshot_path
                else:
                    print(f"Screenshot failed. Return code: {result.returncode}")
                    print(f"STDOUT: {result.stdout}")
                    print(f"STDERR: {result.stderr}")
                    # Clean up
                    if os.path.exists(html_file_path):
                        os.remove(html_file_path)
                    # Fallback to original file
                    return local_file_path
                    
            except subprocess.TimeoutExpired:
                print("Screenshot process timed out, falling back to original file")
                return local_file_path
            
        except Exception as e:
            print(f"Failed to create screenshot for {image_path}: {e}")
            # Return original file as fallback
            return local_file_path
    
    def _prepare_transparent_icon(self, local_file_path: str) -> str:
        """Prepare transparent icon for PPTX to ensure proper transparency and eliminate black backgrounds"""
        try:
            print(f"Processing transparent icon: {local_file_path}")
            
            # Load image
            image = Image.open(local_file_path)
            print(f"Original image mode: {image.mode}, size: {image.size}")
            
            # Always convert to RGBA for consistent processing
            if image.mode != 'RGBA':
                print(f"Converting from {image.mode} to RGBA")
                image = image.convert('RGBA')
            
            # Analyze the image to understand transparency and colors
            width, height = image.size
            pixels = list(image.getdata())
            
            # Count transparency info
            transparent_pixels = sum(1 for r, g, b, a in pixels if a == 0)
            semi_transparent = sum(1 for r, g, b, a in pixels if 0 < a < 255)
            dark_pixels = sum(1 for r, g, b, a in pixels if r < 50 and g < 50 and b < 50 and a > 0)
            total_pixels = len(pixels)
            
            print(f"Transparency analysis: {transparent_pixels} fully transparent, {semi_transparent} semi-transparent, {dark_pixels} dark pixels out of {total_pixels}")
            
            # Enhanced processing for transparent icons to eliminate black backgrounds
            if transparent_pixels > 0 or semi_transparent > 0 or dark_pixels > 0:
                print("Processing icon with enhanced transparency handling...")
                
                # Process each pixel with advanced transparency logic
                new_pixels = []
                for r, g, b, a in pixels:
                    if a == 0:
                        # Fully transparent - keep as is
                        new_pixels.append((0, 0, 0, 0))
                    elif a < 25:
                        # Nearly transparent - make fully transparent to avoid artifacts
                        new_pixels.append((0, 0, 0, 0))
                    elif a < 100:
                        # Semi-transparent - check if it's creating unwanted background
                        if r < 30 and g < 30 and b < 30:
                            # Dark semi-transparent pixel - likely causing black background
                            new_pixels.append((0, 0, 0, 0))
                        else:
                            # Keep with adjusted alpha for cleaner edges
                            new_pixels.append((r, g, b, min(a * 2, 255)))
                    else:
                        # Visible pixel - apply enhanced color processing
                        if r < 15 and g < 15 and b < 15:
                            # Very dark pixel - check if it should be transparent
                            if a < 200:
                                # Dark pixel with medium alpha - make transparent
                                new_pixels.append((0, 0, 0, 0))
                            else:
                                # Keep dark pixel but ensure it's not causing background issues
                                new_pixels.append((max(r, 20), max(g, 20), max(b, 20), min(a, 250)))
                        elif r < 40 and g < 40 and b < 40:
                            # Moderately dark pixel - enhance contrast and clean alpha
                            new_pixels.append((min(r + 15, 255), min(g + 15, 255), min(b + 15, 255), min(a, 250)))
                        else:
                            # Normal colored pixel - keep with clean alpha
                            new_pixels.append((r, g, b, min(a, 250)))
                
                # Create new image with processed pixels
                processed_image = Image.new('RGBA', image.size, (0, 0, 0, 0))
                processed_image.putdata(new_pixels)
                
                # Apply additional smoothing to edges to eliminate artifacts
                from PIL import ImageFilter
                # Light blur to smooth edge artifacts that might cause background issues
                processed_image = processed_image.filter(ImageFilter.GaussianBlur(radius=0.3))
                
                # Enhance the alpha channel to ensure clean transparency
                alpha = processed_image.split()[3]  # Get alpha channel
                # Increase contrast in alpha channel to eliminate gray edge artifacts
                alpha = alpha.point(lambda x: 0 if x < 30 else (255 if x > 225 else min(x * 1.2, 255)))
                
                # Recombine with enhanced alpha
                rgb = processed_image.split()[:3]
                processed_image = Image.merge('RGBA', rgb + (alpha,))
                
                # Save with optimal PNG settings for PPTX transparency
                processed_path = os.path.join(self._temp_dir, f"icon_{uuid.uuid4()}.png")
                processed_image.save(
                    processed_path, 
                    "PNG", 
                    optimize=True,
                    compress_level=6,
                    pnginfo=None  # Remove metadata that might interfere
                )
                
                print(f"Saved enhanced transparent icon to: {processed_path}")
                return processed_path
            else:
                print("Icon has no transparency issues, using original")
                return local_file_path
            
        except Exception as e:
            print(f"Failed to process transparent icon {local_file_path}: {e}")
            # Return original file as fallback
            return local_file_path

    async def fetch_network_assets(self):
        # Skip network asset fetching - python-pptx will handle online URLs directly
        pass

    async def create_ppt(self):
        await self.fetch_network_assets()

        for slide_model in self._slide_models:
            # Adding global shapes to slide
            if self._ppt_model.shapes:
                slide_model.shapes.append(self._ppt_model.shapes)

            self.add_and_populate_slide(slide_model)

    def set_presentation_theme(self):
        slide_master = self._ppt.slide_master
        slide_master_part = slide_master.part

        theme_part = slide_master_part.part_related_by(RT.THEME)
        theme = fromstring(theme_part.blob)

        theme_colors = self._theme.colors.theme_color_mapping
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

        for color_name, hex_value in theme_colors.items():
            if color_name:
                color_element = theme.xpath(
                    f"a:themeElements/a:clrScheme/a:{color_name}/a:srgbClr",
                    namespaces=nsmap,
                )[0]
                color_element.set("val", hex_value.encode("utf-8"))

        theme_part._blob = tostring(theme)

    def add_and_populate_slide(self, slide_model: PptxSlideModel):
        slide = self._ppt.slides.add_slide(self._ppt.slide_layouts[BLANK_SLIDE_LAYOUT])

        if slide_model.background:
            self.apply_fill_to_shape(slide.background, slide_model.background)

        if slide_model.note:
            slide.notes_slide.notes_text_frame.text = slide_model.note

        for shape_model in slide_model.shapes:
            model_type = type(shape_model)

            if model_type is PptxPictureBoxModel:
                self.add_picture(slide, shape_model)

            elif model_type is PptxAutoShapeBoxModel:
                self.add_autoshape(slide, shape_model)

            elif model_type is PptxTextBoxModel:
                self.add_textbox(slide, shape_model)

            elif model_type is PptxConnectorModel:
                self.add_connector(slide, shape_model)

    def add_connector(self, slide: Slide, connector_model: PptxConnectorModel):
        if connector_model.thickness == 0:
            return
        connector_shape = slide.shapes.add_connector(
            connector_model.type, *connector_model.position.to_pt_xyxy()
        )
        connector_shape.line.width = Pt(connector_model.thickness)
        connector_shape.line.color.rgb = RGBColor.from_string(connector_model.color)
        self.set_fill_opacity(connector_shape, connector_model.opacity)

    def add_picture(self, slide: Slide, picture_model: PptxPictureBoxModel):
        image_path = picture_model.picture.path
        
        # Convert URL to local file path
        local_file_path = self._convert_url_to_local_path(image_path)
        
        # Check if we need special processing
        needs_processing = (
            picture_model.border_radius or 
            picture_model.shape == PptxBoxShapeEnum.CIRCLE or
            picture_model.opacity or
            picture_model.invert
        )
        
        try:
            # Check if image is transparent (PNG/SVG icon) 
            is_transparent_icon = self._is_transparent_icon(image_path, local_file_path)
            
            if needs_processing and local_file_path and os.path.exists(local_file_path):
                # Load local image for processing
                image = Image.open(local_file_path)
                image = image.convert("RGBA")
                
                # Apply high-quality processing
                if picture_model.border_radius:
                    image = self._apply_smooth_border_radius(image, picture_model.border_radius)
                
                if picture_model.shape == PptxBoxShapeEnum.CIRCLE:
                    image = self._create_smooth_circle_image(image)
                
                if picture_model.invert:
                    from utils.image_utils import invert_image
                    image = invert_image(image)
                    
                if picture_model.opacity:
                    from utils.image_utils import set_image_opacity
                    image = set_image_opacity(image, picture_model.opacity)
                
                # Save processed image
                processed_path = os.path.join(self._temp_dir, f"{uuid.uuid4()}.png")
                image.save(processed_path, "PNG", quality=95)
                final_path = processed_path
                
            elif is_transparent_icon and local_file_path and os.path.exists(local_file_path):
                # Create background with icon overlay to avoid transparency issues
                print(f"Processing transparent icon: is_transparent_icon={is_transparent_icon}, local_file_path={local_file_path}, file_exists={os.path.exists(local_file_path)}")
                final_path = self._create_icon_with_background(local_file_path, picture_model)
                print(f"After _create_icon_with_background: final_path={final_path}")
                
            elif local_file_path and os.path.exists(local_file_path):
                # Use original local file directly
                final_path = local_file_path
                
            else:
                # Fallback: download from URL if local file not found
                print(f"Local file not found: {local_file_path}, downloading from URL")
                final_path = self._download_image_from_url(image_path, needs_processing, picture_model)
            
            if final_path and os.path.exists(final_path):
                margined_position = self.get_margined_position(
                    picture_model.position, picture_model.margin
                )
                slide.shapes.add_picture(final_path, *margined_position.to_pt_list())
            else:
                print(f"Failed to get valid image path for: {image_path}")
            
        except Exception as e:
            print(f"Failed to process image {image_path}: {e}")
    
    def _convert_url_to_local_path(self, image_path: str) -> str:
        """Convert image URL to local file path"""
        try:
            print(f"Converting URL to local path: {image_path}")
            
            # Handle different URL formats
            if image_path.startswith("http://localhost:5001/static/"):
                # http://localhost:5001/static/icons/xxx.jpg -> /Users/.../servers/fastapi/static/icons/xxx.jpg
                static_path = image_path.replace("http://localhost:5001/static/", "")
                current_file = os.path.abspath(__file__)
                fastapi_dir = os.path.dirname(os.path.dirname(current_file))  # servers/fastapi/
                static_dir = os.path.join(fastapi_dir, "static")
                local_path = os.path.join(static_dir, static_path)
                print(f"Localhost static path: {local_path}")
                return local_path
            elif image_path.startswith("http://localhost:5001/app_data/"):
                # http://localhost:5001/app_data/images/xxx.jpg -> /Users/.../app_data/images/xxx.jpg
                relative_path = image_path.replace("http://localhost:5001/app_data/", "")
                local_path = os.path.join(get_app_data_directory_env(), relative_path)
                print(f"Localhost app data path: {local_path}")
                return local_path
            elif image_path.startswith("https://ppt.samsoncj.xyz/app_data/"):
                # https://ppt.samsoncj.xyz/app_data/images/xxx.jpg -> /Users/.../app_data/images/xxx.jpg
                relative_path = image_path.replace("https://ppt.samsoncj.xyz/app_data/", "")
                local_path = os.path.join(get_app_data_directory_env(), relative_path)
                print(f"App data path: {local_path}")
                return local_path
            elif image_path.startswith("https://ppt.samsoncj.xyz/static/"):
                # https://ppt.samsoncj.xyz/static/icons/xxx.jpg -> /Users/.../servers/fastapi/static/icons/xxx.jpg
                static_path = image_path.replace("https://ppt.samsoncj.xyz/static/", "")
                # Static files are served from fastapi/static directory
                # Current file: /Users/.../servers/fastapi/services/pptx_presentation_creator.py
                # Target:      /Users/.../servers/fastapi/static/
                current_file = os.path.abspath(__file__)
                fastapi_dir = os.path.dirname(os.path.dirname(current_file))  # servers/fastapi/
                static_dir = os.path.join(fastapi_dir, "static")
                local_path = os.path.join(static_dir, static_path)
                print(f"Static path: {local_path}")
                return local_path
            elif image_path.startswith("/app_data/"):
                # /app_data/images/xxx.jpg -> /Users/.../app_data/images/xxx.jpg
                relative_path = image_path.replace("/app_data/", "")
                local_path = os.path.join(get_app_data_directory_env(), relative_path)
                print(f"Relative app data path: {local_path}")
                return local_path
            elif image_path.startswith("/static/"):
                # /static/icons/xxx.jpg -> /Users/.../servers/fastapi/static/icons/xxx.jpg
                static_path = image_path.replace("/static/", "")
                current_file = os.path.abspath(__file__)
                fastapi_dir = os.path.dirname(os.path.dirname(current_file))  # servers/fastapi/
                static_dir = os.path.join(fastapi_dir, "static")
                local_path = os.path.join(static_dir, static_path)
                print(f"Relative static path: {local_path}")
                return local_path
            elif not image_path.startswith("http"):
                # Already a local path or relative path
                if os.path.isabs(image_path):
                    print(f"Absolute path: {image_path}")
                    return image_path
                else:
                    local_path = os.path.join(get_app_data_directory_env(), image_path)
                    print(f"Relative to app_data: {local_path}")
                    return local_path
            else:
                # External URL, no local equivalent
                print(f"External URL, no local equivalent: {image_path}")
                return None
        except Exception as e:
            print(f"Error converting URL to local path: {e}")
            return None
    
    def _download_image_from_url(self, image_path: str, needs_processing: bool, picture_model: PptxPictureBoxModel) -> str:
        """Download image from URL as fallback"""
        try:
            response = requests.get(image_path, timeout=30)
            response.raise_for_status()
            
            if needs_processing:
                from io import BytesIO
                image = Image.open(BytesIO(response.content))
                image = image.convert("RGBA")
                
                # Apply processing (same logic as above)
                if picture_model.border_radius:
                    image = self._apply_smooth_border_radius(image, picture_model.border_radius)
                if picture_model.shape == PptxBoxShapeEnum.CIRCLE:
                    image = self._create_smooth_circle_image(image)
                if picture_model.invert:
                    from utils.image_utils import invert_image
                    image = invert_image(image)
                if picture_model.opacity:
                    from utils.image_utils import set_image_opacity
                    image = set_image_opacity(image, picture_model.opacity)
                
                local_path = os.path.join(self._temp_dir, f"{uuid.uuid4()}.png")
                image.save(local_path, "PNG", quality=95)
            else:
                local_path = os.path.join(self._temp_dir, f"{uuid.uuid4()}.jpg")
                with open(local_path, "wb") as f:
                    f.write(response.content)
            
            return local_path
            
        except Exception as e:
            print(f"Failed to download image from URL {image_path}: {e}")
            return None

    def add_autoshape(self, slide: Slide, autoshape_box_model: PptxAutoShapeBoxModel):
        position = autoshape_box_model.position
        if autoshape_box_model.margin:
            position = self.get_margined_position(position, autoshape_box_model.margin)

        autoshape = slide.shapes.add_shape(
            autoshape_box_model.type, *position.to_pt_list()
        )

        textbox = autoshape.text_frame
        textbox.word_wrap = autoshape_box_model.text_wrap

        self.apply_fill_to_shape(autoshape, autoshape_box_model.fill)
        self.apply_margin_to_text_box(textbox, autoshape_box_model.margin)
        self.apply_stroke_to_shape(autoshape, autoshape_box_model.stroke)
        self.apply_shadow_to_shape(autoshape, autoshape_box_model.shadow)
        self.apply_border_radius_to_shape(autoshape, autoshape_box_model.border_radius)

        if autoshape_box_model.paragraphs:
            self.add_paragraphs(textbox, autoshape_box_model.paragraphs)

    def add_textbox(self, slide: Slide, textbox_model: PptxTextBoxModel):
        position = textbox_model.position
        textbox_shape = slide.shapes.add_textbox(*position.to_pt_list())
        textbox_shape.width += Pt(2)

        textbox = textbox_shape.text_frame
        textbox.word_wrap = textbox_model.text_wrap

        self.apply_fill_to_shape(textbox_shape, textbox_model.fill)
        self.apply_margin_to_text_box(textbox, textbox_model.margin)
        self.add_paragraphs(textbox, textbox_model.paragraphs)

    def add_paragraphs(
        self, textbox: TextFrame, paragraph_models: List[PptxParagraphModel]
    ):
        for index, paragraph_model in enumerate(paragraph_models):
            paragraph = textbox.add_paragraph() if index > 0 else textbox.paragraphs[0]
            self.populate_paragraph(paragraph, paragraph_model)

    def populate_paragraph(
        self, paragraph: _Paragraph, paragraph_model: PptxParagraphModel
    ):
        if paragraph_model.spacing:
            self.apply_spacing_to_paragraph(paragraph, paragraph_model.spacing)

        if paragraph_model.line_height:
            paragraph.line_spacing = paragraph_model.line_height

        if paragraph_model.alignment:
            paragraph.alignment = paragraph_model.alignment

        if paragraph_model.font:
            self.apply_font_to_paragraph(paragraph, paragraph_model.font)

        text_runs = []
        if paragraph_model.text:
            text_runs = self.parse_html_text_to_text_runs(
                paragraph_model.font, paragraph_model.text
            )
        elif paragraph_model.text_runs:
            text_runs = paragraph_model.text_runs

        for text_run_model in text_runs:
            text_run = paragraph.add_run()
            self.populate_text_run(text_run, text_run_model)

    def parse_html_text_to_text_runs(self, font: Optional[PptxFontModel], text: str):
        return parse_inline_html_to_runs(text, font)

    def populate_text_run(self, text_run: _Run, text_run_model: PptxTextRunModel):
        text_run.text = text_run_model.text
        if text_run_model.font:
            self.apply_font(text_run.font, text_run_model.font)

    def apply_border_radius_to_shape(self, shape: Shape, border_radius: Optional[int]):
        if not border_radius:
            return
        try:
            # Simple approach using shape adjustments (for autoshapes only)
            if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
                normalized_border_radius = min(border_radius / min(shape.width.pt, shape.height.pt), 0.5)
                shape.adjustments[0] = normalized_border_radius
                print(f"Applied border radius adjustment: {normalized_border_radius}")
            else:
                print("Shape doesn't support border radius adjustments")
        except Exception as e:
            print(f"Could not apply border radius: {e}")

    def apply_fill_to_shape(self, shape: Shape, fill: Optional[PptxFillModel] = None):
        if not fill:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill.color)
            self.set_fill_opacity(shape.fill, fill.opacity)

    def apply_stroke_to_shape(
        self, shape: Shape, stroke: Optional[PptxStrokeModel] = None
    ):
        if not stroke or stroke.thickness == 0:
            shape.line.fill.background()
        else:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = RGBColor.from_string(stroke.color)
            shape.line.width = Pt(stroke.thickness)
            self.set_fill_opacity(shape.line.fill, stroke.opacity)

    def apply_shadow_to_shape(
        self, shape: Shape, shadow: Optional[PptxShadowModel] = None
    ):

        # Access the XML for the shape
        sp_element = shape._element
        sp_pr = sp_element.xpath("p:spPr")[0]  # Shape properties XML element

        nsmap = sp_pr.nsmap

        # # Remove existing shadow effects if present
        effect_list = sp_pr.find("a:effectLst", namespaces=nsmap)
        if effect_list:
            old_outer_shadow = effect_list.find("a:outerShdw")
            if old_outer_shadow:
                effect_list.remove(
                    old_outer_shadow, namespaces=nsmap
                )  # Remove the old shadow
            old_inner_shadow = effect_list.find("a:innerShdw")
            if old_inner_shadow:
                effect_list.remove(
                    old_inner_shadow, namespaces=nsmap
                )  # Remove the old shadow
            old_prst_shadow = effect_list.find("a:prstShdw")
            if old_prst_shadow:
                effect_list.remove(
                    old_prst_shadow, namespaces=nsmap
                )  # Remove the old shadow

        if not effect_list:
            effect_list = etree.SubElement(
                sp_pr, f"{{{nsmap['a']}}}effectLst", nsmap=nsmap
            )

        if shadow is None:
            # Apply shadow with zero values when shadow is None
            outer_shadow = etree.SubElement(
                effect_list,
                f"{{{nsmap['a']}}}outerShdw",
                {
                    "blurRad": "0",
                    "dist": "0",
                    "dir": "0",
                },
                nsmap=nsmap,
            )
            color_element = etree.SubElement(
                outer_shadow,
                f"{{{nsmap['a']}}}srgbClr",
                {"val": "000000"},
                nsmap=nsmap,
            )
            etree.SubElement(
                color_element,
                f"{{{nsmap['a']}}}alpha",
                {"val": "0"},
                nsmap=nsmap,
            )
        else:
            # Apply the provided shadow
            outer_shadow = etree.SubElement(
                effect_list,
                f"{{{nsmap['a']}}}outerShdw",
                {
                    "blurRad": f"{Pt(shadow.radius)}",
                    "dir": f"{shadow.angle * 1000}",
                    "dist": f"{Pt(shadow.offset)}",
                    "rotWithShape": "0",
                },
                nsmap=nsmap,
            )
            color_element = etree.SubElement(
                outer_shadow,
                f"{{{nsmap['a']}}}srgbClr",
                {"val": f"{shadow.color}"},
                nsmap=nsmap,
            )
            etree.SubElement(
                color_element,
                f"{{{nsmap['a']}}}alpha",
                {"val": f"{int(shadow.opacity * 100000)}"},
                nsmap=nsmap,
            )

    def set_fill_opacity(self, fill, opacity):
        if opacity is None or opacity >= 1.0:
            return

        alpha = int((opacity) * 100000)

        try:
            ts = fill._xPr.solidFill
            sF = ts.get_or_change_to_srgbClr()
            self.get_sub_element(sF, "a:alpha", val=str(alpha))
        except Exception as e:
            print(f"Could not set fill opacity: {e}")

    def get_margined_position(
        self, position: PptxPositionModel, margin: Optional[PptxSpacingModel]
    ) -> PptxPositionModel:
        if not margin:
            return position

        left = position.left + margin.left
        top = position.top + margin.top
        width = max(position.width - margin.left - margin.right, 0)
        height = max(position.height - margin.top - margin.bottom, 0)

        return PptxPositionModel(left=left, top=top, width=width, height=height)

    def apply_margin_to_text_box(
        self, text_frame: TextFrame, margin: Optional[PptxSpacingModel]
    ) -> PptxPositionModel:
        text_frame.margin_left = Pt(margin.left if margin else 0)
        text_frame.margin_right = Pt(margin.right if margin else 0)
        text_frame.margin_top = Pt(margin.top if margin else 0)
        text_frame.margin_bottom = Pt(margin.bottom if margin else 0)

    def apply_spacing_to_paragraph(
        self, paragraph: _Paragraph, spacing: PptxSpacingModel
    ):
        paragraph.space_before = Pt(spacing.top)
        paragraph.space_after = Pt(spacing.bottom)

    def apply_font_to_paragraph(self, paragraph: _Paragraph, font: PptxFontModel):
        self.apply_font(paragraph.font, font)

    def apply_font(self, font: Font, font_model: PptxFontModel):
        font.name = font_model.name
        font.color.rgb = RGBColor.from_string(font_model.color)
        font.italic = font_model.italic
        font.size = Pt(font_model.size)
        font.bold = font_model.font_weight >= 600
        if font_model.underline is not None:
            font.underline = bool(font_model.underline)
        if font_model.strike is not None:
            self.apply_strike_to_font(font, font_model.strike)

    def apply_strike_to_font(self, font: Font, strike: Optional[bool]):
        try:
            rPr = font._element
            if strike is True:
                rPr.set("strike", "sngStrike")
            elif strike is False:
                rPr.set("strike", "noStrike")
        except Exception as e:
            print(f"Could not apply strikethrough: {e}")

    def save(self, path: str):
        self._ppt.save(path)
